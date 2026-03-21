"""
validate_brand.py - Amazon Brand Validator

Checks documents and content for compliance with Amazon brand guidelines.
Validates: colors, fonts, brand name capitalization, and tone keywords.
"""

import re
from dataclasses import dataclass, field
from typing import Optional


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class BrandIssue:
    category: str          # "color" | "font" | "capitalization" | "tone"
    severity: str          # "error" | "warning"
    message: str
    location: Optional[str] = None   # e.g. "Slide 3", "Cell B4", "Page 2"

    def __str__(self) -> str:
        loc = f" [{self.location}]" if self.location else ""
        return f"[{self.severity.upper()}] {self.category}{loc}: {self.message}"


@dataclass
class ValidationReport:
    passed: bool
    issues: list[BrandIssue] = field(default_factory=list)

    @property
    def errors(self) -> list[BrandIssue]:
        return [i for i in self.issues if i.severity == "error"]

    @property
    def warnings(self) -> list[BrandIssue]:
        return [i for i in self.issues if i.severity == "warning"]

    def summary(self) -> str:
        status = "PASS" if self.passed else "FAIL"
        return (
            f"Brand validation: {status} | "
            f"{len(self.errors)} error(s), {len(self.warnings)} warning(s)"
        )

    def __str__(self) -> str:
        lines = [self.summary()]
        for issue in self.issues:
            lines.append(f"  {issue}")
        return "\n".join(lines)


# ---------------------------------------------------------------------------
# Validator
# ---------------------------------------------------------------------------

class BrandValidator:
    """
    Validates content against Amazon brand guidelines.

    Checks:
    - Colors: only approved hex values (#FF9900, #232F3E, #FFFFFF, #000000)
    - Fonts: Segoe UI required; flags Arial/Times New Roman as warnings
    - Brand name: "Amazon" must be capitalized correctly (not "amazon", "AMAZON")
    - Tone: flags off-brand language (e.g., aggressive superlatives, informal slang)

    Usage:
        validator = BrandValidator()

        # Check plain text
        report = validator.validate_text("We are the BEST and most amazin company!")
        print(report)

        # Check a PowerPoint file
        report = validator.validate_powerpoint("deck.pptx")
        print(report)

        # Check an Excel file
        report = validator.validate_excel("report.xlsx")
        print(report)
    """

    APPROVED_COLORS: set[str] = {
        "#FF9900",  # Amazon orange
        "#232F3E",  # Amazon dark navy
        "#FFFFFF",  # White
        "#000000",  # Black
        # Normalised lowercase equivalents
        "#ff9900",
        "#232f3e",
        "#ffffff",
        "#000000",
    }

    APPROVED_FONTS: set[str] = {"Segoe UI", "segoe ui"}
    WARN_FONTS: set[str] = {"Arial", "Helvetica", "Calibri", "Times New Roman"}

    # Brand name patterns
    _BRAND_WRONG_CASE = re.compile(r"\bamazon\b", re.IGNORECASE)
    _BRAND_CORRECT = re.compile(r"\bAmazon\b")
    _BRAND_ALL_CAPS = re.compile(r"\bAMAZON\b")

    # Tone keywords
    TONE_ERRORS: list[str] = [
        "cheapest", "dirt cheap", "crap", "sucks", "hate", "useless",
    ]
    TONE_WARNINGS: list[str] = [
        "best ever", "greatest of all time", "unbeatable", "unstoppable",
        "insane deal", "crazy price", "literally", "awesome sauce",
    ]

    # ------------------------------------------------------------------ #
    #  Public API
    # ------------------------------------------------------------------ #

    def validate_text(self, text: str, location: Optional[str] = None) -> ValidationReport:
        """Validate a plain-text string for brand compliance."""
        issues: list[BrandIssue] = []
        issues.extend(self._check_brand_name(text, location))
        issues.extend(self._check_tone(text, location))
        return self._make_report(issues)

    def validate_colors(
        self, hex_colors: list[str], location: Optional[str] = None
    ) -> ValidationReport:
        """Validate a list of hex color strings against approved brand colors."""
        issues: list[BrandIssue] = []
        for color in hex_colors:
            normalised = color.upper() if color.startswith("#") else f"#{color.upper()}"
            if normalised not in {c.upper() for c in self.APPROVED_COLORS}:
                issues.append(BrandIssue(
                    category="color",
                    severity="error",
                    message=f"Unapproved color '{color}'. Approved: #FF9900, #232F3E, #FFFFFF, #000000",
                    location=location,
                ))
        return self._make_report(issues)

    def validate_fonts(
        self, font_names: list[str], location: Optional[str] = None
    ) -> ValidationReport:
        """Validate font names against the brand font standard (Segoe UI)."""
        issues: list[BrandIssue] = []
        for font in font_names:
            if font.lower() not in {f.lower() for f in self.APPROVED_FONTS}:
                severity = "warning" if font in self.WARN_FONTS else "error"
                issues.append(BrandIssue(
                    category="font",
                    severity=severity,
                    message=f"Non-brand font '{font}'. Required: Segoe UI",
                    location=location,
                ))
        return self._make_report(issues)

    def validate_powerpoint(self, file_path: str) -> ValidationReport:
        """
        Full brand validation of a PowerPoint file.

        Requires: python-pptx  (pip install python-pptx)
        """
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError("python-pptx is required: pip install python-pptx")

        prs = Presentation(file_path)
        issues: list[BrandIssue] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            loc = f"Slide {slide_num}"
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                full_text = " ".join(
                    run.text
                    for para in shape.text_frame.paragraphs
                    for run in para.runs
                )
                issues.extend(self._check_brand_name(full_text, loc))
                issues.extend(self._check_tone(full_text, loc))

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Font check
                        if run.font.name:
                            sub = self.validate_fonts([run.font.name], loc)
                            issues.extend(sub.issues)
                        # Color check
                        try:
                            if run.font.color and run.font.color.type:
                                rgb: RGBColor = run.font.color.rgb
                                hex_color = f"#{rgb}"
                                sub = self.validate_colors([hex_color], loc)
                                issues.extend(sub.issues)
                        except Exception:
                            pass

        return self._make_report(issues)

    def validate_excel(self, file_path: str) -> ValidationReport:
        """
        Full brand validation of an Excel workbook.

        Requires: openpyxl  (pip install openpyxl)
        """
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path)
        issues: list[BrandIssue] = []

        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    loc = f"Sheet '{ws.title}' cell {cell.coordinate}"

                    # Text checks
                    if cell.value and isinstance(cell.value, str):
                        issues.extend(self._check_brand_name(cell.value, loc))
                        issues.extend(self._check_tone(cell.value, loc))

                    # Font checks
                    if cell.font:
                        if cell.font.name:
                            sub = self.validate_fonts([cell.font.name], loc)
                            issues.extend(sub.issues)
                        if cell.font.color and cell.font.color.rgb:
                            hex_color = f"#{cell.font.color.rgb}"
                            sub = self.validate_colors([hex_color], loc)
                            issues.extend(sub.issues)

                    # Fill color checks
                    if cell.fill and cell.fill.fgColor and cell.fill.fgColor.rgb:
                        hex_color = f"#{cell.fill.fgColor.rgb}"
                        if hex_color.upper() not in ("#FF000000", "#FF9900FF"):  # skip transparent
                            sub = self.validate_colors([hex_color], loc)
                            issues.extend(sub.issues)

        return self._make_report(issues)

    # ------------------------------------------------------------------ #
    #  Private helpers
    # ------------------------------------------------------------------ #

    def _check_brand_name(
        self, text: str, location: Optional[str]
    ) -> list[BrandIssue]:
        issues = []
        # Find all occurrences of the word (any casing)
        for match in self._BRAND_WRONG_CASE.finditer(text):
            word = match.group()
            if self._BRAND_CORRECT.match(word):
                continue  # correct capitalisation
            if self._BRAND_ALL_CAPS.match(word):
                issues.append(BrandIssue(
                    category="capitalization",
                    severity="error",
                    message=f"Brand name written as '{word}'. Must be 'Amazon' (title case only).",
                    location=location,
                ))
            else:
                issues.append(BrandIssue(
                    category="capitalization",
                    severity="error",
                    message=f"Brand name written as '{word}'. Must be 'Amazon'.",
                    location=location,
                ))
        return issues

    def _check_tone(self, text: str, location: Optional[str]) -> list[BrandIssue]:
        issues = []
        lower = text.lower()
        for kw in self.TONE_ERRORS:
            if kw in lower:
                issues.append(BrandIssue(
                    category="tone",
                    severity="error",
                    message=f"Off-brand language detected: '{kw}'.",
                    location=location,
                ))
        for kw in self.TONE_WARNINGS:
            if kw in lower:
                issues.append(BrandIssue(
                    category="tone",
                    severity="warning",
                    message=f"Potentially off-brand phrase: '{kw}'. Review for brand alignment.",
                    location=location,
                ))
        return issues

    @staticmethod
    def _make_report(issues: list[BrandIssue]) -> ValidationReport:
        passed = not any(i.severity == "error" for i in issues)
        return ValidationReport(passed=passed, issues=issues)


# ---------------------------------------------------------------------------
# CLI convenience
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python validate_brand.py <file.pptx|file.xlsx|'text to check'>")
        sys.exit(1)

    target = sys.argv[1]
    validator = BrandValidator()

    if target.endswith(".pptx"):
        report = validator.validate_powerpoint(target)
    elif target.endswith(".xlsx"):
        report = validator.validate_excel(target)
    else:
        report = validator.validate_text(target)

    print(report)
    sys.exit(0 if report.passed else 1)
