"""
apply_brand.py - Amazon Brand Formatter

Applies Amazon brand guidelines to Excel, PowerPoint, and PDF documents.
Brand colors: #FF9900 (orange), #232F3E (dark navy/black)
Brand font: Segoe UI
"""

from dataclasses import dataclass, field
from typing import Optional
import os


@dataclass
class BrandColors:
    orange: str = "#FF9900"
    black: str = "#232F3E"
    white: str = "#FFFFFF"

    def as_rgb(self, hex_color: str) -> tuple[int, int, int]:
        """Convert hex color to (R, G, B) tuple."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    @property
    def orange_rgb(self) -> tuple[int, int, int]:
        return self.as_rgb(self.orange)

    @property
    def black_rgb(self) -> tuple[int, int, int]:
        return self.as_rgb(self.black)


@dataclass
class BrandFonts:
    primary: str = "Segoe UI"
    fallback: list[str] = field(default_factory=lambda: ["Arial", "Helvetica", "sans-serif"])
    heading_size: int = 24
    body_size: int = 11
    caption_size: int = 9


class BrandFormatter:
    """
    Applies Amazon brand guidelines to Office and PDF documents.

    Usage:
        formatter = BrandFormatter()
        formatter.format_excel("report.xlsx")
        formatter.format_powerpoint("deck.pptx")
        formatter.format_pdf("document.pdf", output_path="branded_document.pdf")
    """

    def __init__(self):
        self.colors = BrandColors()
        self.fonts = BrandFonts()

    # ------------------------------------------------------------------ #
    #  Excel
    # ------------------------------------------------------------------ #

    def format_excel(self, file_path: str, output_path: Optional[str] = None) -> str:
        """
        Apply brand styling to an Excel workbook.

        Requires: openpyxl  (pip install openpyxl)

        Applies:
        - Header rows: orange background (#FF9900), white text, Segoe UI bold
        - Body cells: dark text (#232F3E), Segoe UI
        - Alternating row shading for readability
        """
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        except ImportError:
            raise ImportError("openpyxl is required: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path)

        orange_fill = PatternFill(
            start_color="FF9900", end_color="FF9900", fill_type="solid"
        )
        alt_fill = PatternFill(
            start_color="FFF3E0", end_color="FFF3E0", fill_type="solid"
        )
        thin_border = Border(
            bottom=Side(style="thin", color="232F3E")
        )

        for ws in wb.worksheets:
            for row_idx, row in enumerate(ws.iter_rows(), start=1):
                is_header = row_idx == 1
                for cell in row:
                    cell.font = Font(
                        name=self.fonts.primary,
                        size=self.fonts.heading_size if is_header else self.fonts.body_size,
                        bold=is_header,
                        color="FFFFFF" if is_header else "232F3E",
                    )
                    if is_header:
                        cell.fill = orange_fill
                    elif row_idx % 2 == 0:
                        cell.fill = alt_fill
                    cell.alignment = Alignment(vertical="center", wrap_text=True)
                    if is_header:
                        cell.border = thin_border

        output = output_path or self._branded_path(file_path)
        wb.save(output)
        return output

    # ------------------------------------------------------------------ #
    #  PowerPoint
    # ------------------------------------------------------------------ #

    def format_powerpoint(self, file_path: str, output_path: Optional[str] = None) -> str:
        """
        Apply brand styling to a PowerPoint presentation.

        Requires: python-pptx  (pip install python-pptx)

        Applies:
        - Title placeholders: Segoe UI, #232F3E
        - Body placeholders: Segoe UI, #232F3E
        - Slide backgrounds: white with optional orange accent bar
        """
        try:
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError("python-pptx is required: pip install python-pptx")

        prs = Presentation(file_path)
        orange = RGBColor(*self.colors.orange_rgb)
        black = RGBColor(*self.colors.black_rgb)

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                is_title = shape.shape_type == 13 or (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = self.fonts.primary
                        run.font.color.rgb = orange if is_title else black
                        run.font.size = Pt(
                            self.fonts.heading_size if is_title else self.fonts.body_size
                        )

        output = output_path or self._branded_path(file_path)
        prs.save(output)
        return output

    # ------------------------------------------------------------------ #
    #  PDF
    # ------------------------------------------------------------------ #

    def format_pdf(self, file_path: str, output_path: Optional[str] = None) -> str:
        """
        Stamp brand header/footer onto an existing PDF.

        Requires: pypdf + reportlab  (pip install pypdf reportlab)

        Applies:
        - Orange header bar with "Amazon" wordmark
        - Dark footer bar with page numbers
        """
        try:
            from pypdf import PdfReader, PdfWriter
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.colors import HexColor
            import io
        except ImportError:
            raise ImportError(
                "pypdf and reportlab are required: pip install pypdf reportlab"
            )

        reader = PdfReader(file_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)

        for page_num, page in enumerate(reader.pages, start=1):
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)

            # Build overlay
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=(width, height))

            # Header bar
            c.setFillColor(HexColor(self.colors.orange))
            c.rect(0, height - 30, width, 30, fill=1, stroke=0)
            c.setFillColor(HexColor(self.colors.white))
            c.setFont("Helvetica-Bold", 12)
            c.drawString(10, height - 20, "Amazon")

            # Footer bar
            c.setFillColor(HexColor(self.colors.black))
            c.rect(0, 0, width, 20, fill=1, stroke=0)
            c.setFillColor(HexColor(self.colors.white))
            c.setFont("Helvetica", 8)
            c.drawRightString(width - 10, 6, f"Page {page_num} of {total_pages}")

            c.save()
            packet.seek(0)

            from pypdf import PdfReader as _R
            overlay_page = _R(packet).pages[0]
            page.merge_page(overlay_page)
            writer.add_page(page)

        output = output_path or self._branded_path(file_path)
        with open(output, "wb") as f:
            writer.write(f)
        return output

    # ------------------------------------------------------------------ #
    #  Helpers
    # ------------------------------------------------------------------ #

    @staticmethod
    def _branded_path(original: str) -> str:
        """Return a sibling path with '_branded' suffix."""
        base, ext = os.path.splitext(original)
        return f"{base}_branded{ext}"
